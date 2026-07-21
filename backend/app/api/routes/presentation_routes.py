from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape
import traceback
from app.controllers.presentation_controller import PresentationController
from app.schemas.presentation_schema import (
    PresentationCreate,
    PresentationUpdate,
)

router = APIRouter(
    prefix="/api/presentations",
    tags=["Presentations"],
)


from fastapi import UploadFile, File
import shutil
import os
import uuid

@router.post("/upload")
async def upload_media(file: UploadFile = File(...)):
    try:
        # Create unique filename
        ext = file.filename.split(".")[-1]
        filename = f"{uuid.uuid4()}.{ext}"
        filepath = os.path.join("uploads", filename)
        
        with open(filepath, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
            
        return {"url": f"http://127.0.0.1:8000/uploads/{filename}"}
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/")
async def create_presentation(presentation: PresentationCreate):
    try:
        return await PresentationController.create(
            presentation.model_dump()
        )
    except Exception as e:
        traceback.print_exc()
        return {"error": str(e)}


@router.get("/")
async def get_presentations():
    return await PresentationController.get_all()


@router.get("/{presentation_id}")
async def get_presentation(
    presentation_id: str,
):
    return await PresentationController.get_by_id(
        presentation_id
    )


@router.put("/{presentation_id}")
async def update_presentation(
    presentation_id: str,
    presentation: PresentationUpdate,
):
    return await PresentationController.update(
        presentation_id,
        presentation.model_dump(exclude_unset=True),
    )

@router.delete("/{presentation_id}")
async def delete_presentation(
    presentation_id: str,
):
    return await PresentationController.delete(
        presentation_id
    )

@router.get("/{presentation_id}/export/pptx")
async def export_presentation_pptx(
    presentation_id: str,
):
    presentation_data = (
        await PresentationController.get_by_id(
            presentation_id
        )
    )

    if not presentation_data:
        raise HTTPException(
            status_code=404,
            detail="Presentation not found",
        )

    # Create PowerPoint presentation
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = presentation_data.get(
        "slides",
        [],
    )

    for slide_data in slides:
        # Add blank slide
        slide = prs.slides.add_slide(
            prs.slide_layouts[6]
        )

        # Get Fabric canvas JSON
        canvas_json = slide_data.get(
            "canvasJSON",
            {},
        )

        if not canvas_json:
            canvas_json = {}

        # Get Fabric objects
        objects = canvas_json.get(
            "objects",
            [],
        )

        for obj in objects:
            obj_type = obj.get("type")

            print(
                "EXPORT OBJECT:",
                obj_type,
            )
            prs.slide_layouts[6]
        

        # Get Fabric canvas JSON
        canvas_json = slide_data.get(
            "canvasJSON",
            {},
        )

        if not canvas_json:
            canvas_json = {}

        # Get Fabric objects
        objects = canvas_json.get(
            "objects",
            [],
        )

        for obj in objects:
            obj_type = obj.get("type")

            print(
                "EXPORT OBJECT:",
                obj_type,
            )

            # Export Fabric text objects
            if obj_type in [
                "text",
                "textbox",
                "i-text",
            ]:
                # Fabric usually uses left/top
                x = obj.get(
                    "left",
                    obj.get("x", 100),
                )

                y = obj.get(
                    "top",
                    obj.get("y", 100),
                )

                width = obj.get(
                    "width",
                    300,
                )

                height = obj.get(
                    "height",
                    100,
                )

                # Include Fabric scaling
                scale_x = obj.get(
                    "scaleX",
                    1,
                )

                scale_y = obj.get(
                    "scaleY",
                    1,
                )

                width = width * scale_x
                height = height * scale_y

                # Pixels -> Inches
                left = Inches(x / 96)
                top = Inches(y / 96)

                box_width = Inches(
                    max(width, 10) / 96
                )

                box_height = Inches(
                    max(height, 10) / 96
                )

                textbox = (
                    slide.shapes.add_textbox(
                        left,
                        top,
                        box_width,
                        box_height,
                    )
                )

                text_frame = (
                    textbox.text_frame
                )

                text_frame.text = obj.get(
                    "text",
                    obj.get(
                        "content",
                        "",
                    ),
                )

                font_size = obj.get(
                    "fontSize",
                    24,
                )

                for paragraph in (
                    text_frame.paragraphs
                ):
                    for run in (
                        paragraph.runs
                    ):
                        run.font.size = Pt(
                            font_size
                        )

    # If there are no slides,
    # create one blank slide
    if len(prs.slides) == 0:
        prs.slides.add_slide(
            prs.slide_layouts[6]
        )

    # Save PPTX in memory
    output = BytesIO()

    prs.save(output)

    output.seek(0)

    # File name
    title = presentation_data.get(
        "title",
        "presentation",
    )

    safe_title = (
        title
        .replace(" ", "_")
        .replace("/", "_")
        .replace("\\", "_")
    )

    # Return PPTX file
    return StreamingResponse(
        output,
        media_type=(
            "application/vnd.openxmlformats-"
            "officedocument.presentationml."
            "presentation"
        ),
        headers={
            "Content-Disposition": (
                f'attachment; '
                f'filename="{safe_title}.pptx"'
            )
        },
    )

import os
import glob
import subprocess
import asyncio
from playwright.async_api import async_playwright

@router.get("/{presentation_id}/export/pdf")
async def export_presentation_pdf(presentation_id: str):
    presentation_data = await PresentationController.get_by_id(presentation_id)
    if not presentation_data:
        raise HTTPException(status_code=404, detail="Presentation not found")

    title = presentation_data.get("title", "presentation")
    safe_title = title.replace(" ", "_").replace("/", "_").replace("\\", "_")
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page()
        url = f"http://localhost:5173/?presentationId={presentation_id}&mode=pdf"
        await page.goto(url, wait_until="networkidle")
        
        await page.wait_for_timeout(2000)
        
        pdf_bytes = await page.pdf(
            format="Letter",
            landscape=True,
            print_background=True
        )
        await browser.close()
        
    return StreamingResponse(
        BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_title}.pdf"'
        },
    )

@router.get("/{presentation_id}/export/mp4")
async def export_presentation_mp4(presentation_id: str):
    presentation_data = await PresentationController.get_by_id(presentation_id)
    if not presentation_data:
        raise HTTPException(status_code=404, detail="Presentation not found")

    title = presentation_data.get("title", "presentation")
    safe_title = title.replace(" ", "_").replace("/", "_").replace("\\", "_")
    
    video_dir = f"./temp_video_{presentation_id}"
    os.makedirs(video_dir, exist_ok=True)
    
    slides = presentation_data.get("slides", [])
    total_time_ms = len(slides) * 5000 + 2000 
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            record_video_dir=video_dir,
            record_video_size={"width": 1280, "height": 720}
        )
        page = await context.new_page()
        url = f"http://localhost:5173/?presentationId={presentation_id}&mode=mp4"
        
        await page.goto(url)
        await page.wait_for_timeout(total_time_ms)
        
        await context.close()
        await browser.close()
        
    webm_files = glob.glob(f"{video_dir}/*.webm")
    if not webm_files:
        raise HTTPException(status_code=500, detail="Video recording failed")
        
    webm_file = webm_files[0]
    mp4_file = f"{video_dir}/output.mp4"
    
    subprocess.run([
        "ffmpeg", "-y", "-i", webm_file, "-c:v", "libx264", "-c:a", "aac", mp4_file
    ], check=True)
    
    with open(mp4_file, "rb") as f:
        mp4_bytes = f.read()
        
    os.remove(webm_file)
    os.remove(mp4_file)
    os.rmdir(video_dir)
    
    return StreamingResponse(
        BytesIO(mp4_bytes),
        media_type="video/mp4",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_title}.mp4"'
        },
    )
